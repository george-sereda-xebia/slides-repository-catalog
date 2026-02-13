"""PowerPoint slides renderer using LibreOffice headless."""

import os
import logging
import signal
import subprocess
import shutil
import tempfile
import uuid
from pathlib import Path
from typing import List, Dict

from pptx import Presentation


logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class SlidesRenderer:
    """Renderer for converting PPTX to PDF and PNG images."""

    def __init__(self, output_dir: str):
        """Initialize slides renderer.

        Args:
            output_dir: Directory to store rendered output
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def find_libreoffice(self) -> str:
        """Find LibreOffice executable path.

        Returns:
            Path to soffice executable

        Raises:
            RuntimeError: If LibreOffice not found
        """
        paths = [
            "soffice",
            "/usr/bin/soffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",
            "C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe",
        ]

        for path in paths:
            if shutil.which(path) or os.path.exists(path):
                logger.info(f"Found LibreOffice at: {path}")
                return path

        raise RuntimeError(
            "LibreOffice not found. Please install LibreOffice:\n"
            "  macOS: brew install libreoffice\n"
            "  Ubuntu: sudo apt-get install libreoffice\n"
            "  Windows: Download from https://www.libreoffice.org/"
        )

    def _kill_stale_soffice(self):
        """Kill any lingering soffice processes to avoid profile locking."""
        try:
            subprocess.run(
                ["pkill", "-f", "soffice"],
                capture_output=True,
                timeout=5,
            )
        except Exception:
            pass

    def convert_to_pdf(self, pptx_path: str, output_dir: str) -> str:
        """Convert PPTX to PDF using LibreOffice with a temporary user profile.

        Uses a unique temp profile per conversion to avoid profile locking issues.

        Args:
            pptx_path: Path to PPTX file
            output_dir: Directory to write the PDF into

        Returns:
            Path to generated PDF file

        Raises:
            RuntimeError: If conversion fails
        """
        self._kill_stale_soffice()

        soffice = self.find_libreoffice()
        profile_dir = tempfile.mkdtemp(prefix="lo_profile_")

        try:
            cmd = [
                soffice,
                f"-env:UserInstallation=file://{profile_dir}",
                "--headless",
                "--norestore",
                "--nofirststartwizard",
                "--convert-to", "pdf",
                "--outdir", output_dir,
                pptx_path,
            ]

            logger.info(f"Converting to PDF: {Path(pptx_path).name}")
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=60,
            )

            if result.returncode != 0:
                logger.error(f"LibreOffice stderr: {result.stderr}")
                raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")

            # Find the generated PDF
            pptx_stem = Path(pptx_path).stem
            pdf_path = Path(output_dir) / f"{pptx_stem}.pdf"

            if not pdf_path.exists():
                raise RuntimeError(f"PDF not created at expected path: {pdf_path}")

            logger.info(f"PDF created: {pdf_path}")
            return str(pdf_path)

        except subprocess.TimeoutExpired:
            logger.error("LibreOffice conversion timed out (60s)")
            self._kill_stale_soffice()
            raise RuntimeError("LibreOffice conversion timed out")
        finally:
            shutil.rmtree(profile_dir, ignore_errors=True)

    def get_slide_count(self, pptx_path: str) -> int:
        """Get the number of slides in a PPTX file.

        Args:
            pptx_path: Path to PPTX file

        Returns:
            Number of slides
        """
        try:
            prs = Presentation(pptx_path)
            return len(prs.slides)
        except Exception as e:
            logger.warning(f"Failed to count slides in {pptx_path}: {e}")
            return 0

    def render_pptx(self, pptx_path: str, presentation_id: str) -> List[str]:
        """Render PPTX slides to PNG images.

        Args:
            pptx_path: Path to PPTX file
            presentation_id: Unique ID for this presentation

        Returns:
            List of paths to rendered PNG files
        """
        logger.info(f"Rendering PNGs: {pptx_path}")

        presentation_dir = self.output_dir / presentation_id
        presentation_dir.mkdir(parents=True, exist_ok=True)

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            soffice = self.find_libreoffice()

            cmd = [
                soffice,
                "--headless",
                "--convert-to", "png",
                "--outdir", str(temp_path),
                pptx_path,
            ]

            try:
                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=300,
                    check=True,
                )
                logger.debug(f"LibreOffice output: {result.stdout}")
            except subprocess.CalledProcessError as e:
                logger.error(f"LibreOffice conversion failed: {e.stderr}")
                raise
            except subprocess.TimeoutExpired:
                logger.error("LibreOffice conversion timed out (5 minutes)")
                raise

            temp_files = sorted(temp_path.glob("*.png"))

            if not temp_files:
                logger.warning(f"No PNG files generated for {pptx_path}")
                return []

            output_files = []
            for idx, temp_file in enumerate(temp_files, start=1):
                output_file = presentation_dir / f"slide_{idx:03d}.png"
                shutil.move(str(temp_file), str(output_file))
                output_files.append(str(output_file))

            logger.info(f"Rendered {len(output_files)} slides for {presentation_id}")
            return output_files

    def extract_text(self, pptx_path: str) -> str:
        """Extract all text from PPTX file.

        Args:
            pptx_path: Path to PPTX file

        Returns:
            Extracted text as string
        """
        try:
            prs = Presentation(pptx_path)
            text_parts = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            text_parts.append(text)

            extracted_text = " ".join(text_parts)
            logger.debug(f"Extracted {len(extracted_text)} characters from {Path(pptx_path).name}")
            return extracted_text

        except Exception as e:
            logger.warning(f"Failed to extract text from {pptx_path}: {e}")
            return ""

    def render_presentation(self, pptx_path: str, presentation_id: str) -> dict:
        """Convert presentation to PDF and return metadata.

        Args:
            pptx_path: Path to PPTX file
            presentation_id: Unique ID for this presentation

        Returns:
            Dictionary with pdf_path, slide_count, and text
        """
        try:
            # Create output directory for this presentation's PDF
            presentation_dir = self.output_dir / presentation_id
            presentation_dir.mkdir(parents=True, exist_ok=True)

            # Convert to native PDF
            pdf_path = self.convert_to_pdf(pptx_path, str(presentation_dir))

            # Get slide count
            slide_count = self.get_slide_count(pptx_path)

            # Extract text for search metadata
            text_content = self.extract_text(pptx_path)

            return {
                "success": True,
                "pdf_path": pdf_path,
                "slide_count": slide_count,
                "text": text_content,
                "error": None,
            }

        except Exception as e:
            logger.error(f"Failed to render {pptx_path}: {e}")
            return {
                "success": False,
                "pdf_path": None,
                "slide_count": 0,
                "text": "",
                "error": str(e),
            }


def main():
    """Test the renderer."""
    import sys

    if len(sys.argv) < 2:
        print("Usage: python slides_renderer.py <path-to-pptx>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    renderer = SlidesRenderer("./output")
    presentation_id = Path(pptx_path).stem
    result = renderer.render_presentation(pptx_path, presentation_id)

    if result["success"]:
        print(f"\nSuccess! PDF: {result['pdf_path']}, Slides: {result['slide_count']}")
    else:
        print(f"\nError: {result['error']}")
        sys.exit(1)


if __name__ == "__main__":
    main()
